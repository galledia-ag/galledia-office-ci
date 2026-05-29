"""
helpers.py — Galledia Präsentations-Skill v1.3 (adaptive Headlines + adaptive Karten)
Schrift-Regel: Volte (Fliesstext/Kapiteltitel), Volte Semibold (Headlines/Hervorhebungen)
Keine Versalien. Kein Volte Rounded in generiertem Inhalt.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ── CI-Farb-Tokens ────────────────────────────────────────────────────────────
RED  = RGBColor(0xE6,0x1C,0x52); BLACK = RGBColor(0,0,0); WHITE = RGBColor(0xFF,0xFF,0xFF)
G1   = RGBColor(0x40,0x40,0x40); G2    = RGBColor(0x66,0x66,0x66)
G3   = RGBColor(0xA6,0xA6,0xA6); G4    = RGBColor(0xD9,0xD9,0xD9)
GL   = RGBColor(0xF2,0xF2,0xF5); TURK  = RGBColor(0x22,0xAA,0x9F)

# ── CI-Schrift-Tokens ─────────────────────────────────────────────────────────
BODY = "Volte"           # Fliesstext, Kapiteltitel, Labels
SB   = "Volte Semibold"  # Headlines, Kartenüberschriften, Hervorhebungen

# ── Template-Grid (aus Vorlage_5.pptx, idx=0/11 auf 04_vielText) ──────────────
GX   = Inches(0.979)   # Linke Hilfslinie
GY_K = Inches(0.750)   # Y Kapiteltitel
GY_H = Inches(1.061)   # Y Headline
GW   = Inches(18.026)  # Content-Breite
SH   = Inches(11.25)   # Folienhöhe

LAYOUTS = {
    "titel":"Titelfolie","abschluss":"Abschlussfolie","schluss":"Schlussfolie",
    "zwischenrot":"2_Zwischenfolie rot","agenda5":"01_Agenda 5","agenda22":"01_Agenda 22",
    "wenig":"02_wenigText","viel":"04_vielText","leer":"Leer","piktogramme":"Piktogramme",
}
_TEMPLATE = os.path.join(os.path.dirname(__file__),"assets","Vorlage_5.pptx")
_LOGO_DIR  = os.path.join(os.path.dirname(__file__),"assets","logo")
LOGOS = {  # CI: Rot→helle Folien | Weiss→rote/dunkle Folien | Schwarz→fehlt (neu exportieren)
    "rot":              os.path.join(_LOGO_DIR,"logo_rot.png"),
    "rot_schriftzug":   os.path.join(_LOGO_DIR,"logo_rot_schriftzug.png"),
    "weiss":            os.path.join(_LOGO_DIR,"logo_weiss.png"),
    "weiss_schriftzug":   os.path.join(_LOGO_DIR,"logo_weiss_schriftzug.png"),
    "schwarz":            os.path.join(_LOGO_DIR,"logo_schwarz.png"),
    "schwarz_schriftzug": os.path.join(_LOGO_DIR,"logo_schwarz_schriftzug.png"),
}

# ── Basis-Setup ───────────────────────────────────────────────────────────────

def build_presentation(template=None, datum="", rechtseinheit=""):
    """
    Vorlage laden, Beispielfolien entfernen.
    datum / rechtseinheit werden als Fusszeilen-Defaults gespeichert.
    Beispiel: build_presentation(datum='29. Mai 2026', rechtseinheit='Galledia Fachmedien AG')
    """
    prs = Presentation(template or _TEMPLATE)
    sl = prs.slides._sldIdLst
    for sid in list(sl):
        prs.part.drop_rel(sid.get(qn("r:id"))); sl.remove(sid)
    prs._layouts = {l.name:l for l in prs.slide_masters[0].slide_layouts}
    prs._datum = datum
    prs._rechtseinheit = rechtseinheit
    _strip_transitions(prs)
    return prs


def _strip_transitions(prs):
    """Entfernt alle Folienübergänge aus Master, Layouts und Slides."""
    TAG = qn('p:transition')
    sources = []
    for m in prs.slide_masters:
        sources.append(m._element)
        for lay in m.slide_layouts:
            sources.append(lay._element)
    for slide in prs.slides:
        sources.append(slide._element)
    for el in sources:
        for tr in el.findall('.//' + TAG):
            tr.getparent().remove(tr)


def _move_placeholder(ph, x_emu, y_emu, w_emu, h_emu):
    """Setzt Position+Grösse eines Platzhalters explizit (überschreibt Layout-Vererbung)."""
    from lxml import etree
    sp = ph._element
    sp_pr = sp.find(qn('p:spPr'))
    if sp_pr is None:
        sp_pr = etree.SubElement(sp, qn('p:spPr'))
    xfrm = sp_pr.find(qn('a:xfrm'))
    if xfrm is None:
        xfrm = etree.SubElement(sp_pr, qn('a:xfrm'))
    off = xfrm.find(qn('a:off'))
    if off is None:
        off = etree.SubElement(xfrm, qn('a:off'))
    off.set('x', str(int(x_emu))); off.set('y', str(int(y_emu)))
    ext = xfrm.find(qn('a:ext'))
    if ext is None:
        ext = etree.SubElement(xfrm, qn('a:ext'))
    ext.set('cx', str(int(w_emu))); ext.set('cy', str(int(h_emu)))


def _lay(prs, key):
    name = LAYOUTS.get(key, key)
    cache = getattr(prs,"_layouts",{}) or {l.name:l for l in prs.slide_masters[0].slide_layouts}
    if name not in cache:
        raise ValueError(f"Layout '{name}' fehlt. Verfügbar: {list(cache)}")
    return cache[name]

def _blank(prs): return prs.slides.add_slide(_lay(prs,"leer"))

def _phs(slide): return {p.placeholder_format.idx:p for p in slide.placeholders}

# ── Primitive ─────────────────────────────────────────────────────────────────

def tb(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    b=slide.shapes.add_textbox(x,y,w,h); tf=b.text_frame
    tf.word_wrap=True; tf.vertical_anchor=anchor
    tf.margin_left=tf.margin_right=tf.margin_top=tf.margin_bottom=0
    return tf

def run(para, text, size, font, color,
        align=PP_ALIGN.LEFT, space_after=6, bold=False, italic=False):
    para.text=text; para.alignment=align; para.space_after=Pt(space_after)
    r=para.runs[0]; r.font.name=font; r.font.size=Pt(size)
    r.font.color.rgb=color; r.font.bold=bold; r.font.italic=italic
    return para

def _headline_size(text):
    """Adaptive Headline-Grösse — verhindert Überlauf in den Body bei langen Headlines."""
    n = len(text or "")
    if n <= 28: return 72
    if n <= 38: return 54
    if n <= 50: return 44
    if n <= 64: return 36
    return 30

def _enable_shrink(tf):
    """Aktiviert 'Text bei Überlauf verkleinern' (normAutofit) auf einem Textframe."""
    from lxml import etree
    bodyPr = tf._txBody.find(qn('a:bodyPr'))
    if bodyPr is None:
        bodyPr = etree.SubElement(tf._txBody, qn('a:bodyPr'))
    for tag in ('a:spAutoFit','a:noAutofit','a:normAutofit'):
        e = bodyPr.find(qn(tag))
        if e is not None: bodyPr.remove(e)
    etree.SubElement(bodyPr, qn('a:normAutofit'))

def card(slide, x, y, w, h, fill, corner=0.06):
    sh=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,x,y,w,h)
    sh.fill.solid(); sh.fill.fore_color.rgb=fill
    sh.line.fill.background(); sh.shadow.inherit=False
    try: sh.adjustments[0]=corner
    except: pass
    return sh

# ── Interne Header/Footer-Helfer ──────────────────────────────────────────────

def _footer_str(prs, folio):
    """Baut den Fusszeilen-Text: 'Folie [Nr], [Datum], [Rechtseinheit]'"""
    parts = []
    if folio: parts.append(f"Folie {folio}")
    if getattr(prs,"_datum",""): parts.append(prs._datum)
    if getattr(prs,"_rechtseinheit",""): parts.append(prs._rechtseinheit)
    return ", ".join(parts)

def _set_header(slide, kicker, headline):
    """
    Setzt Kapiteltitel + Headline auf exakter Template-Position.
    Füllt idx=0 und idx=11 direkt (kein Versatz), idx=13 wird geleert.
    Kapiteltitel: 30 pt Volte, Schwarz. Headline: 72 pt Volte Semibold, Schwarz.
    """
    ph = _phs(slide)
    if 0 in ph:
        run(ph[0].text_frame.paragraphs[0], kicker, 30, BODY, BLACK, space_after=0)
    else:
        tf=tb(slide,GX,GY_K,GW,Inches(0.35))
        run(tf.paragraphs[0], kicker, 30, BODY, BLACK, space_after=0)
    if headline:
        if 11 in ph:
            run(ph[11].text_frame.paragraphs[0], headline, _headline_size(headline), SB, BLACK, space_after=0); _enable_shrink(ph[11].text_frame)
        else:
            tf=tb(slide,GX,GY_H,GW,Inches(1.15))
            run(tf.paragraphs[0], headline, _headline_size(headline), SB, BLACK, space_after=0)
    if 13 in ph and not ph[13].text_frame.text:
        ph[13].text=""   # Prompt-Text unterdrücken

def _set_footer(slide, prs, folio, source=""):
    """Fusszeile links: 'Folie [Nr], [Datum], [Rechtseinheit]'. Quellenangabe darüber."""
    ph = _phs(slide)
    footer = _footer_str(prs, folio)
    if footer:
        if 14 in ph:
            run(ph[14].text_frame.paragraphs[0], footer, 11, BODY, BLACK,
                align=PP_ALIGN.LEFT, space_after=0)
        else:
            # Textbox auf exakter Hilfslinie (y=10.82", h=0.25", w=GW — wie idx=14 auf 04_vielText)
            tf=tb(slide,GX,Inches(10.82),GW,Inches(0.25))
            run(tf.paragraphs[0], footer, 11, BODY, BLACK, space_after=0)
    if source:
        tf=tb(slide,GX,SH-Inches(1.2),GW,Inches(0.35))
        run(tf.paragraphs[0],f"Quelle: {source}",11,BODY,G2,italic=True,space_after=0)

# ── Native-Layout-Slides ──────────────────────────────────────────────────────

# Titelfolie-Geometrie (aus Vorlage_5 gemessen)
_TITLE_Y_SUB        = Inches(3.968)   # y idx=10 Untertitel (bei 2-zeiligem Haupttitel)
_TITLE_Y_MAIN       = Inches(4.683)   # y idx=0  Haupttitel
_TITLE_1LINE_SHIFT  = Inches(0.68)    # Verschiebung nach unten bei 1-zeiligem Haupttitel
_TITLE_CHARS_1LINE  = 40              # Schwellenwert: ≤ 1 Zeile; > 2 Zeilen

def add_title(prs, title, subtitle=""):
    """
    Titelfolie. Dynamische Zentrierung:
    Einzeiliger Haupttitel (≤40 Zeichen) → beide Platzhalter um 0.68" nach unten,
    sodass der Block optisch zentriert bleibt.
    """
    s=prs.slides.add_slide(_lay(prs,"titel"))
    ph=_phs(s); ph[0].text=title
    if 10 in ph and subtitle: ph[10].text=subtitle
    if len(title) <= _TITLE_CHARS_1LINE:
        shift = int(_TITLE_1LINE_SHIFT)
        # Gemessene Layout-Koordinaten (aus Vorlage_5): x, y, w, h in EMU
        specs = {
            0:  (Inches(1.023), Inches(4.683), Inches(17.955), Inches(2.175)),
            10: (Inches(1.023), Inches(3.968), Inches(17.955), Inches(0.752)),
        }
        for idx in (0, 10):
            if idx in ph:
                x, y, w, h = specs[idx]
                _move_placeholder(ph[idx], x, int(y) + shift, w, h)
    return s

def add_section(prs, number, title):
    """Optionaler Kapitel-Anker (für Decks >15 Folien mit mehreren Kapiteln)."""
    s=prs.slides.add_slide(_lay(prs,"zwischenrot"))
    tf=tb(s,GX,Inches(3.8),GW,Inches(4),MSO_ANCHOR.TOP)
    run(tf.paragraphs[0], number, 54, SB, WHITE, space_after=4)
    run(tf.add_paragraph(), title,  40, SB, WHITE)
    return s

def add_agenda(prs, items, variant="agenda5", folio=""):
    """
    Agenda-Folie. «Agenda» steht in Kapiteltitel (idx=0) UND Haupttitel (idx=11).
    variant: 'agenda5' (≤5 Punkte) oder 'agenda22' (6–12 Punkte).
    """
    s=prs.slides.add_slide(_lay(prs,variant)); ph=_phs(s)
    ph[0].text="Agenda"; ph[11].text="Agenda"
    tf=ph[13].text_frame; tf.paragraphs[0].text=items[0]
    for item in items[1:]: tf.add_paragraph().text=item
    if 14 in ph: run(ph[14].text_frame.paragraphs[0],
                     _footer_str(prs,folio),11,BODY,BLACK,space_after=0)
    return s

def add_content(prs, variant, kapitel, headline, body_text, folio="", source=""):
    """Inhaltsfolie. variant='viel' oder 'wenig'."""
    s=prs.slides.add_slide(_lay(prs,variant)); ph=_phs(s)
    ph[0].text=kapitel
    run(ph[11].text_frame.paragraphs[0], headline, _headline_size(headline), SB, BLACK, space_after=0)
    _enable_shrink(ph[11].text_frame)
    ph[13].text=body_text
    _enable_shrink(ph[13].text_frame)
    if 14 in ph: run(ph[14].text_frame.paragraphs[0],
                     _footer_str(prs,folio),11,BODY,BLACK,space_after=0)
    if source and 15 in ph: ph[15].text=f"Quelle: {source}"
    return s

def add_closing(prs):   return prs.slides.add_slide(_lay(prs,"schluss"))
def add_discussion(prs): return prs.slides.add_slide(_lay(prs,"abschluss"))

def add_logo(slide, variant="rot", size_in=1.2):
    """Bildmarke rechts unten. variant: 'rot' (helle Folien) | 'weiss' (rote/dunkle)."""
    from PIL import Image as _Im
    path = LOGOS.get(variant, LOGOS["rot"])
    w_px, h_px = _Im.open(path).size
    w_emu = Inches(size_in); h_emu = int(w_emu * h_px / w_px)
    slide.shapes.add_picture(path, int(GX+GW-w_emu), int(SH-Inches(0.5)-h_emu), w_emu, h_emu)


# ── Kompositions-Komponenten (Leer-Layout) ────────────────────────────────────

def kpi_grid(prs, kpis, kicker="", headline="", folio="", source=""):
    s=_blank(prs)
    if kicker or headline: _set_header(s,kicker,headline)
    n=len(kpis); g=Inches(0.5); cw=(GW-(n-1)*g)/n; cy,ch=Inches(3.6),Inches(4.0)
    for i,(zahl,label) in enumerate(kpis):
        x=GX+i*(cw+g); fill=RED if i==0 else GL
        fg=WHITE if i==0 else BLACK; fg2=WHITE if i==0 else G2
        card(s,x,cy,cw,ch,fill)
        tf=tb(s,x+Inches(0.3),cy,cw-Inches(0.6),ch,MSO_ANCHOR.MIDDLE)
        run(tf.paragraphs[0],zahl, 56,SB,  fg, align=PP_ALIGN.CENTER,space_after=8)
        run(tf.add_paragraph(),label,17,BODY,fg2,align=PP_ALIGN.CENTER,space_after=0)
    _set_footer(s,prs,folio,source); return s

def two_column(prs,head_l,items_l,head_r,items_r,
               col2_red=True,kicker="",headline="",folio="",source=""):
    s=_blank(prs)
    if kicker or headline: _set_header(s,kicker,headline)
    # Kartenhöhe an Inhalt anpassen (Header + Bullets), beide Karten gleich hoch
    rows=max(len(items_l),len(items_r))
    ch=Inches(min(6.8, max(2.4, 1.5 + 0.62*rows)))
    cy=Inches(3.2); cw=(GW-Inches(0.6))/2
    for x,fill,head,fg,items in [
        (GX, GL, head_l, G1, items_l),
        (GX+cw+Inches(0.6), RED if col2_red else GL,
         head_r, WHITE if col2_red else G1, items_r)]:
        card(s,x,cy,cw,ch,fill)
        tf=tb(s,x+Inches(0.6),cy+Inches(0.5),cw-Inches(1.2),ch-Inches(1.0))
        run(tf.paragraphs[0],head,28,SB,  fg,space_after=16)
        for item in items: run(tf.add_paragraph(),f"•  {item}",19,BODY,fg,space_after=10)
    _set_footer(s,prs,folio,source); return s

def flow_pipeline(prs,nodes,cache_text="",kicker="",headline="",folio="",source=""):
    s=_blank(prs)
    if kicker or headline: _set_header(s,kicker,headline)
    n=len(nodes); g=Inches(0.4); nw=(GW-(n-1)*g)/n; ny,nh=Inches(3.6),Inches(2.6)
    for i,node in enumerate(nodes):
        titel=node[0] if isinstance(node,(list,tuple)) else node
        sub  =node[1] if isinstance(node,(list,tuple)) and len(node)>1 else ""
        fill=RED if i==n-1 else GL; fg=WHITE if i==n-1 else BLACK; fg2=WHITE if i==n-1 else G2
        x=GX+i*(nw+g); card(s,x,ny,nw,nh,fill)
        tf=tb(s,x+Inches(0.2),ny,nw-Inches(0.4),nh,MSO_ANCHOR.MIDDLE)
        run(tf.paragraphs[0],titel,20,SB,  fg, align=PP_ALIGN.CENTER,space_after=6)
        if sub: run(tf.add_paragraph(),sub,13,BODY,fg2,align=PP_ALIGN.CENTER,space_after=0)
        if i<n-1:
            at=tb(s,x+nw,ny,g,nh,MSO_ANCHOR.MIDDLE)
            run(at.paragraphs[0],"→",22,SB,G2,align=PP_ALIGN.CENTER,space_after=0)
    if cache_text:
        cy2=ny+nh+Inches(0.55); card(s,GX,cy2,GW,Inches(1.4),TURK)
        tf=tb(s,GX+Inches(0.6),cy2,GW-Inches(1.2),Inches(1.4),MSO_ANCHOR.MIDDLE)
        run(tf.paragraphs[0],"Semantic Caching",18,SB,  WHITE,space_after=4)
        run(tf.add_paragraph(),cache_text,      15,BODY,WHITE,space_after=0)
    _set_footer(s,prs,folio,source); return s

def timeline(prs,phases,kicker="",headline="",folio="",source=""):
    s=_blank(prs)
    if kicker or headline: _set_header(s,kicker,headline)
    ly=Inches(5.4); line=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,GX,ly,GW,Pt(4))
    line.fill.solid(); line.fill.fore_color.rgb=G4; line.line.fill.background(); line.shadow.inherit=False
    n=len(phases); seg=GW/n
    for i,(titel,beschr) in enumerate(phases):
        cx=GX+seg*i+seg/2; col=RED if i==0 else G2
        dot=s.shapes.add_shape(MSO_SHAPE.OVAL,int(cx-Inches(0.22)),int(ly-Inches(0.18)),Inches(0.44),Inches(0.44))
        dot.fill.solid(); dot.fill.fore_color.rgb=col; dot.line.color.rgb=WHITE; dot.line.width=Pt(3); dot.shadow.inherit=False
        tf=tb(s,int(cx-seg/2+Inches(0.3)),ly-Inches(2.6),int(seg-Inches(0.6)),Inches(2.3),MSO_ANCHOR.BOTTOM)
        run(tf.paragraphs[0],titel,22,SB,  BLACK,align=PP_ALIGN.CENTER,space_after=0)
        tf2=tb(s,int(cx-seg/2+Inches(0.3)),ly+Inches(0.5),int(seg-Inches(0.6)),Inches(2.0))
        run(tf2.paragraphs[0],beschr,16,BODY,G2,align=PP_ALIGN.CENTER,space_after=0)
    _set_footer(s,prs,folio,source); return s

def numbered_steps(prs,steps,kicker="",headline="",folio="",source=""):
    s=_blank(prs)
    if kicker or headline: _set_header(s,kicker,headline)
    n=len(steps); g=Inches(0.5); cw=(GW-(n-1)*g)/n; cy,ch=Inches(3.6),Inches(5.0)
    for i,(titel,beschr) in enumerate(steps):
        x=GX+i*(cw+g); card(s,x,cy,cw,ch,GL)
        badge=s.shapes.add_shape(MSO_SHAPE.OVAL,x+Inches(0.6),cy+Inches(0.55),Inches(1.0),Inches(1.0))
        badge.fill.solid(); badge.fill.fore_color.rgb=RED; badge.line.fill.background(); badge.shadow.inherit=False
        bf=badge.text_frame; bf.word_wrap=False
        run(bf.paragraphs[0],str(i+1),32,SB,WHITE,align=PP_ALIGN.CENTER,space_after=0)
        tf=tb(s,x+Inches(0.6),cy+Inches(1.9),cw-Inches(1.2),ch-Inches(2.3))
        run(tf.paragraphs[0],titel, 24,SB,  BLACK,space_after=10)
        run(tf.add_paragraph(),beschr,17,BODY,G2,  space_after=0)
    _set_footer(s,prs,folio,source); return s

def image_bleed(prs,image_path,kicker="",headline="",body_text="",
                folio="",source="",image_side="right"):
    s=_blank(prs); img_w=Inches(9.2)
    img_x=Inches(20)-img_w if image_side=="right" else 0
    txt_x=GX if image_side=="right" else img_w+Inches(0.5)
    s.shapes.add_picture(image_path,img_x,0,img_w,SH)
    if kicker or headline: _set_header(s,kicker,headline)
    if body_text:
        tf=tb(s,txt_x,Inches(2.4),Inches(9.0),Inches(6.0))
        run(tf.paragraphs[0],body_text,19,BODY,G1,space_after=0)
    _set_footer(s,prs,folio,source); return s
