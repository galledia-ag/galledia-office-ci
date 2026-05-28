"""
fill_praesentation.py — generiert eine Galledia-Praesentation aus der .potx-Vorlage.

V1-Scope: erste Folie (Titelseite) mit Titel + Untertitel befuellen,
MasterProperty-Platzhalter ersetzen, optional weitere Slides hinzufuegen,
fertige .pptx schreiben.

Aufruf:
    python fill_praesentation.py --input data.json --output deck.pptx
    python fill_praesentation.py --input - --output deck.pptx   (stdin)
"""

from __future__ import annotations

import argparse
import io
import json
import re
import sys
import zipfile
from copy import deepcopy
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt

SCRIPT_DIR = Path(__file__).resolve().parent
SKILL_DIR = SCRIPT_DIR.parent
TEMPLATE_PATH = SKILL_DIR / "templates" / "Präsentationsvorlage Galledia.potx"

TEMPLATE_CT = "application/vnd.openxmlformats-officedocument.presentationml.template.main+xml"
PRESENTATION_CT = "application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml"

VALID_OE = {
    "galledia group ag",
    "Galledia Fachmedien AG",
    "Galledia Regionalmedien AG",
    "Galledia Print AG",
    "Galledia Digital AG",
}

FORBIDDEN = [
    (re.compile(r"\bGalledia AG\b"), "Es gibt keine 'Galledia AG' — exakte OE verwenden"),
    (re.compile(r"\bGalledia Gruppe\b"), "Heisst 'galledia group ag' (klein geschrieben)"),
    (re.compile(r"\bGalledia GmbH\b"), "Alle Galledia-OE sind AG, keine GmbH"),
    (re.compile(r"(?i)\bfax\b"), "Fax wird nicht mehr verwendet"),
]

# Layout-Mapping (Index → Layout-Name in der Vorlage). Wir referenzieren
# Layouts via Schluesselwort, damit der User nicht den exakten Namen kennen muss.
LAYOUT_BY_KEY = {
    "title":      "Titelfolie rot",
    "title_red":  "Titelfolie rot",
    "title_turquoise": "Titelfolie türkis",
    "title_bronze":    "Titelfolie bronze",
    "title_purple":    "Titelfolie purple",
    "title_lila":      "Titelfolie lila",
    "section_red":  "Zwischenfolie rot",
    "section":      "Zwischenfolie rot",
    "section_turquoise": "Zwischenfolie türkis",
    "section_bronze":    "Zwischenfolie bronze",
    "section_purple":    "Zwischenfolie purple",
    "section_lila":      "Zwischenfolie lila",
    "agenda":     "01_Agenda 5",
    "agenda_22":  "01_Agenda 22",
    "content":    "02_wenigText",
    "content_long": "04_vielText",
    "default":    "DEFAULT SLIDE",
    "blank":      "Leer",
}


class ValidationError(Exception):
    pass


# ---------------------------------------------------------------------------
# Input decoding / mojibake repair (identisch zu Brief/Kurzbrief)
# ---------------------------------------------------------------------------

MOJIBAKE_MARKERS = (
    "Ã¼", "Ã¶", "Ã¤", "ÃŸ", "Ãœ", "Ã–", "Ã„", "Ã©", "Ã¨", "Ã ", "Ã§",
    "Â·", "Â«", "Â»", "Â§", "Â°", "â€™", "â€œ", "â€",
)


def _decode_input(raw: bytes) -> str:
    if raw.startswith(b"\xef\xbb\xbf"):
        return raw[3:].decode("utf-8")
    try:
        return raw.decode("utf-8")
    except UnicodeDecodeError:
        return raw.decode("cp1252")


def _repair_mojibake(obj):
    if isinstance(obj, dict):
        return {k: _repair_mojibake(v) for k, v in obj.items()}
    if isinstance(obj, list):
        return [_repair_mojibake(v) for v in obj]
    if isinstance(obj, str):
        return _fix_mojibake_string(obj)
    return obj


def _fix_mojibake_string(s: str) -> str:
    if not s or not any(m in s for m in MOJIBAKE_MARKERS):
        return s
    try:
        return s.encode("latin-1").decode("utf-8")
    except (UnicodeEncodeError, UnicodeDecodeError):
        return s


# ---------------------------------------------------------------------------
# Validation
# ---------------------------------------------------------------------------

def validate(data: dict) -> None:
    errors: list[str] = []
    oe = data.get("organisation", "")
    if oe and oe not in VALID_OE:
        errors.append(
            f"organisation '{oe}' ist keine gueltige OE. "
            f"Erlaubt: {sorted(VALID_OE)}"
        )
    body_text = " ".join(_collect_all_text(data))
    for pat, msg in FORBIDDEN:
        if pat.search(body_text):
            errors.append(f"Verbotene Schreibweise: {msg}")
    if errors:
        raise ValidationError("\n".join(f"- {e}" for e in errors))


def _collect_all_text(obj):
    if isinstance(obj, str):
        yield obj
    elif isinstance(obj, dict):
        for v in obj.values():
            yield from _collect_all_text(v)
    elif isinstance(obj, list):
        for v in obj:
            yield from _collect_all_text(v)


# ---------------------------------------------------------------------------
# .potx -> .pptx (Content-Type swap, sonst lehnt python-pptx ab)
# ---------------------------------------------------------------------------

def _potx_to_pptx_bytes(template_path: Path) -> bytes:
    with zipfile.ZipFile(str(template_path), "r") as zin:
        buf = io.BytesIO()
        with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zout:
            for item in zin.infolist():
                data = zin.read(item.filename)
                if item.filename == "[Content_Types].xml":
                    data = data.replace(TEMPLATE_CT.encode(), PRESENTATION_CT.encode())
                zout.writestr(item, data)
        return buf.getvalue()


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

# Officeatwork MasterProperty-Platzhalter im Vorlagentext, z.B.
# [[MasterProperty("Organisation", "Organisation")]] -> "Galledia Fachmedien AG"
MASTERPROP_RE = re.compile(r'\[\[MasterProperty\("([^"]+)",\s*"([^"]+)"\)\]\]')


def _replace_masterproperty(text: str, organisation: str) -> str:
    """Ersetzt alle MasterProperty-Platzhalter im Text durch konkrete Werte."""
    def repl(match):
        key = match.group(1)
        if key == "Organisation":
            return organisation
        return ""  # unbekannte Property: leeren
    return MASTERPROP_RE.sub(repl, text)


def _replace_text_in_textframe(tf, organisation: str) -> None:
    """Ersetzt MasterProperty-Platzhalter in einem TextFrame; preserves runs."""
    for para in tf.paragraphs:
        for run in para.runs:
            if "[[MasterProperty" in run.text:
                run.text = _replace_masterproperty(run.text, organisation)


def _set_textframe_text(tf, lines: list[str]) -> None:
    """
    Setzt den Text eines TextFrames neu, behaelt aber das Format des ersten Runs
    im ersten Paragraph. Mehrere Lines -> mehrere Paragraphs.
    """
    if not lines:
        lines = [""]
    # Behalte das Format des ersten Run als Template
    if tf.paragraphs and tf.paragraphs[0].runs:
        template_run = tf.paragraphs[0].runs[0]
        font_name = template_run.font.name
        font_size = template_run.font.size
        font_bold = template_run.font.bold
        font_color = None
        try:
            if template_run.font.color and template_run.font.color.rgb:
                font_color = template_run.font.color.rgb
        except (AttributeError, ValueError):
            pass
    else:
        font_name = font_size = font_bold = font_color = None

    tf.clear()
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = line
        if font_name:
            run.font.name = font_name
        if font_size:
            run.font.size = font_size
        if font_bold is not None:
            run.font.bold = font_bold
        if font_color is not None:
            try:
                run.font.color.rgb = font_color
            except (AttributeError, ValueError):
                pass


def _find_layout(prs: Presentation, key_or_name: str):
    """Findet ein Slide Layout nach Schluessel (LAYOUT_BY_KEY) oder direktem Namen."""
    target_name = LAYOUT_BY_KEY.get(key_or_name, key_or_name)
    for layout in prs.slide_layouts:
        if layout.name == target_name:
            return layout
    # Fallback: Substring-Match
    for layout in prs.slide_layouts:
        if target_name.lower() in layout.name.lower():
            return layout
    return prs.slide_layouts[0]  # absoluter Fallback


def _remove_all_slides(prs: Presentation) -> None:
    """Loescht alle bestehenden Slides aus der Praesentation."""
    sldIdLst = prs.slides._sldIdLst
    for sldId in list(sldIdLst):
        rId = sldId.get(
            "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
        )
        prs.part.drop_rel(rId)
        sldIdLst.remove(sldId)


def _add_slide(prs: Presentation, slide_def: dict, organisation: str):
    """
    Fuegt eine Slide nach Definition hinzu.
    slide_def: {layout: str, title: str, subtitle/content: str|list}
    """
    layout = _find_layout(prs, slide_def.get("layout", "content"))
    slide = prs.slides.add_slide(layout)

    title = slide_def.get("title", "")
    subtitle = slide_def.get("subtitle", "")
    content = slide_def.get("content", slide_def.get("items", ""))

    # Title-Placeholder finden (idx=0 oder type=TITLE)
    title_ph = None
    body_phs: list = []
    for shape in slide.placeholders:
        idx = shape.placeholder_format.idx
        if idx == 0:
            title_ph = shape
        else:
            body_phs.append(shape)

    if title_ph and title:
        _set_textframe_text(
            title_ph.text_frame,
            [_replace_masterproperty(title, organisation)]
        )

    # Untertitel/Content in den ersten Body-Placeholder
    if body_phs:
        body_text = subtitle if subtitle else content
        if isinstance(body_text, list):
            lines = [_replace_masterproperty(str(s), organisation) for s in body_text]
        else:
            lines = [_replace_masterproperty(str(body_text), organisation)]
            # split inner \n in single string
            if len(lines) == 1 and "\n" in lines[0]:
                lines = lines[0].split("\n")
        if any(l for l in lines):
            _set_textframe_text(body_phs[0].text_frame, lines)

    return slide


# ---------------------------------------------------------------------------
# Main fill
# ---------------------------------------------------------------------------

def fill(data: dict, template: Path, output: Path) -> dict:
    validate(data)
    organisation = data.get("organisation", "Galledia Fachmedien AG")

    report: dict = {}

    # Open template
    pptx_bytes = _potx_to_pptx_bytes(template)
    prs = Presentation(io.BytesIO(pptx_bytes))

    slides_def = data.get("slides", [])
    replace_existing = data.get("replace_existing_slides", True)

    if slides_def:
        if replace_existing:
            _remove_all_slides(prs)
        for slide_def in slides_def:
            _add_slide(prs, slide_def, organisation)
        report["slides_added"] = len(slides_def)
    else:
        # Keine slides angegeben: bestehende Vorlage-Slides behalten,
        # nur MasterProperty global ersetzen.
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    _replace_text_in_textframe(shape.text_frame, organisation)
        report["slides_added"] = 0
        report["existing_slides_kept"] = len(prs.slides)

    # Globaler MasterProperty-Replace auf allen Slides (auch neu hinzugefuegte)
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                _replace_text_in_textframe(shape.text_frame, organisation)

    # Save
    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output))
    report["output"] = str(output)
    report["organisation"] = organisation
    return report


def main() -> int:
    ap = argparse.ArgumentParser(description="Galledia Praesentation Generator")
    ap.add_argument("--input", required=True, help="Pfad zur JSON-Eingabe oder '-' fuer stdin")
    ap.add_argument("--output", required=True, help="Pfad fuer Ausgabe-.pptx")
    ap.add_argument("--template", default=str(TEMPLATE_PATH), help="Pfad zur .potx-Vorlage")
    args = ap.parse_args()

    if args.input == "-":
        raw = sys.stdin.buffer.read()
        data = json.loads(_decode_input(raw))
    else:
        with open(args.input, "rb") as f:
            data = json.loads(_decode_input(f.read()))

    data = _repair_mojibake(data)

    try:
        report = fill(data, Path(args.template), Path(args.output))
    except ValidationError as e:
        print("VALIDIERUNG FEHLGESCHLAGEN:", file=sys.stderr)
        print(str(e), file=sys.stderr)
        return 2

    print(json.dumps(report, indent=2, ensure_ascii=False))
    return 0


if __name__ == "__main__":
    sys.exit(main())
